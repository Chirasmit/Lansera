import os
import cv2
import numpy as np
import pytesseract
import tensorflow as tf
from tensorflow.keras.applications import MobileNetV2
from tensorflow.keras.applications.mobilenet_v2 import preprocess_input, decode_predictions
import cohere
import fitz  
from pptx import Presentation  
from docx import Document  
import tempfile
import streamlit as st

def convert_pptx_to_text(file):
    prs = Presentation(file)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def convert_pdf_to_text(file):
    try:
        doc = fitz.open(stream=file.read(), filetype="pdf")
        text = []
        for page in doc:
            text.append(page.get_text("text")) 
        return "\n".join(text)
    except Exception as e:
        st.error(f"Failed to process PDF: {e}")
        return None

def convert_docx_to_text(file):
    doc = Document(file)
    text = []
    for para in doc.paragraphs:
        text.append(para.text)
    return "\n".join(text)


def get_file_text(file):
    file_ext = os.path.splitext(file.name)[1].lower()

    if file_ext == ".txt":
        return file.read().decode("utf-8")
    elif file_ext == ".pptx":
        return convert_pptx_to_text(file)
    elif file_ext == ".docx":
        return convert_docx_to_text(file)
    elif file_ext == ".pdf":
        return convert_pdf_to_text(file)
    else:
        st.error("Unsupported file format. Please use a .txt, .pptx, .docx, or .pdf file.")
        return None

def initialize_cohere(api_key):
    return cohere.Client(api_key)

def query_cohere(client, question, context):
    try:
        response = client.chat(
             message=question,
            documents=[{"text": context}],
            temperature=0.7,
        )
        return response.text
    except Exception as e:
        return f"error while querying"


def load_and_preprocess_image(image_path):
    image = cv2.imread(image_path)
    if image is None:
        raise FileNotFoundError(f"Image not found at {image_path}. Please check the path.")
    image = cv2.cvtColor(image, cv2.COLOR_BGR2RGB)
    image = cv2.resize(image, (224, 224))
    return preprocess_input(image)

def predict_image(image_path):
    model = MobileNetV2(weights='imagenet')
    preprocessed_image = load_and_preprocess_image(image_path)
    preprocessed_image = np.expand_dims(preprocessed_image, axis=0)
    predictions = model.predict(preprocessed_image)
    decoded_predictions = decode_predictions(predictions, top=3)[0]

    result = []
    for _, class_name, confidence in decoded_predictions:
        result.append(f"{class_name}: {confidence * 100:.2f}%")
    return result

def extract_images_from_pdf(file, output_dir):
    doc = fitz.open(stream=file.read(), filetype="pdf")
    images = []
    for i, page in enumerate(doc):
        for img_index, img in enumerate(page.get_images(full=True)):
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            image_ext = base_image["ext"]
            image_path = os.path.join(output_dir, f"page_{i + 1}img{img_index + 1}.{image_ext}")
            with open(image_path, "wb") as img_file:
                img_file.write(image_bytes)
            images.append(image_path)
    return images

def extract_images_from_docx(file, output_dir):
    doc = Document(file)
    images = []
    for rel in doc.part.rels.values():
        if "image" in rel.target_ref:
            image = rel.target_part.blob
            image_path = os.path.join(output_dir, rel.target_ref.split("/")[-1])
            with open(image_path, "wb") as img_file:
                img_file.write(image)
            images.append(image_path)
    return images

def extract_images_from_pptx(file, output_dir):
    prs = Presentation(file)
    images = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13: 
                image = shape.image.blob
                image_ext = shape.image.ext
                image_path = os.path.join(output_dir, f"slide_{prs.slides.index(slide) + 1}.{image_ext}")
                with open(image_path, "wb") as img_file:
                    img_file.write(image)
                images.append(image_path)
    return images


st.title("Lansera")

option = st.sidebar.selectbox("Choose an option:", [
    "Analyze Document and Ask Questions",
    "Predict Content of an Image",
    "Analyze Document and Find Relevant Images"
])


COHERE_API_KEY = os.getenv('bP4LYKuyDiF6D4G2VQl7Bg4oYtZG2jgYpz2CtMKD', 'bP4LYKuyDiF6D4G2VQl7Bg4oYtZG2jgYpz2CtMKD') 

if not COHERE_API_KEY:
    st.error("API key is missing! Please set it in your environment variables or Streamlit secrets.")
else:
    cohere_client = initialize_cohere(COHERE_API_KEY)

if option == "Analyze Document and Ask Questions":
    uploaded_file = st.file_uploader("Upload a document file:", type=["txt", "pdf", "pptx", "docx"])

    if uploaded_file:
        file_text = get_file_text(uploaded_file)

        if file_text:
            st.success("File loaded successfully! Now you can ask questions about the document.")

            question = st.text_input("Enter your question:")
            if question:
                answer = query_cohere(cohere_client, question, file_text)
                st.write("*Answer:*")
                st.write(answer)


elif option == "Predict Content of an Image":
    uploaded_image = st.file_uploader("Upload an image file:", type=["jpg", "jpeg", "png"])

    if uploaded_image:
        with tempfile.NamedTemporaryFile(delete=False) as temp_image:
            temp_image.write(uploaded_image.read())
            predictions = predict_image(temp_image.name)
            st.image(uploaded_image, caption="Uploaded Image", use_column_width=True)
            st.write("*Predictions:*")
            for prediction in predictions:
                st.write(prediction)

elif option == "Analyze Document and Find Relevant Images":
    api_key = st.text_input("Enter your Cohere API key:", type="password")
    uploaded_file = st.file_uploader("Upload a document file:", type=["txt", "pdf", "pptx", "docx"])

    if uploaded_file and api_key:
        file_text = get_file_text(uploaded_file)

        if file_text:
            st.success("File loaded successfully! Now you can ask questions about the document.")

            question = st.text_input("Enter your question:")
            if question:
                cohere_client = initialize_cohere(api_key)
                answer = query_cohere(cohere_client, question, file_text)
                st.write("*Answer from the document:*")
                st.write(answer)

                with tempfile.TemporaryDirectory() as temp_dir:
                    if uploaded_file.name.lower().endswith(".pdf"):
                        images = extract_images_from_pdf(uploaded_file, temp_dir)
                    elif uploaded_file.name.lower().endswith(".docx"):
                        images = extract_images_from_docx(uploaded_file, temp_dir)
                    elif uploaded_file.name.lower().endswith(".pptx"):
                        images = extract_images_from_pptx(uploaded_file, temp_dir)
                    else:
                        images = []

                    if images:
                        st.write("*Extracted Images:*")
                        for img_path in images:
                            st.image(img_path, use_column_width=True)

                            image_text = pytesseract.image_to_string(img_path)
                            if any(word in image_text.lower() for word in question.lower().split()):
                                st.write(f"*Image matches question context:* {img_path}")